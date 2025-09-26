import io
import os
from tempfile import TemporaryDirectory

import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from supplementary.readers.reader_interface import INPUT_TYPE, FormatReader
from supplementary.utils.data_fetching import (
    get_extension,
    get_material_name,
    get_pmcid,
)
from utils.logging.logging_setup import supplementary_error_logger
from supplementary.utils.result_saver import get_output_type, save


MAX_FILE_SIZE = 50000000  # 50 MB


class PPTXFormatReader(FormatReader):
    def read_bytes(
        self,
        byte_contents,
        source: str,
        type=INPUT_TYPE.BYTES,
    ):
        if not byte_contents:
            return
        pptx_file = io.BytesIO(byte_contents)
        # Check the file size
        file_size = len(byte_contents)
        if file_size > MAX_FILE_SIZE:
            error_message = f"FILE TO LARGE -> {file_size=:_}\t {MAX_FILE_SIZE=:_}"
            supplementary_error_logger.error("%s | %s", source, error_message)
            return None

        file_name = "_".join(source.split("/")[-2:])  # PMCid_materialName
        if (extension := get_extension(source)) == "ppt":
            with TemporaryDirectory() as tempdir:
                ppt_file = f"{tempdir}/{file_name}"
                file_name_no_extension = ppt_file[: -len(extension)]
                pptx_file = f"{tempdir}/{file_name_no_extension}.pptx"
                with open(ppt_file, "wb") as f:
                    f.write(byte_contents)
                try:
                    os.system(
                        f"cd {tempdir} ; libreoffice --headless --convert-to pptx {file_name}"
                    )
                    with open(pptx_file, "rb") as f:
                        pptx_file = io.BytesIO(f.read())
                except Exception as e:
                    # TODO: Implement error handling
                    supplementary_error_logger.error(
                        "%s | %s", source, str(e), exc_info=True
                    )
                    return None
        return self.read(pptx_file, type, source)

    # def read_from_path(self, path, type=INPUT_TYPE.PATH):
    #     return self.read(path, type)
    #     pass

    def read(self, pptx_file, type: INPUT_TYPE, source: str = ""):
        try:
            presentation = Presentation(pptx_file)
        except KeyError as e:
            """
            happened when trying to open a powerpoint inside the zip file
                KeyError: "no relationship of type 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument' in collection"
            """
            supplementary_error_logger.error("%s | %s", source, str(e), exc_info=True)
            return None
        except Exception as e:
            supplementary_error_logger.error("%s | %s", source, str(e), exc_info=True)
            return None
        # to_save = print_out_pptx(presentation)
        to_save = process_pptx(presentation, source)
        if self.save:
            save(
                get_pmcid(source),
                get_material_name(source),
                to_save,
                get_output_type(get_extension(source)),
            )
        return to_save


def process_pptx(presentation: Presentation, source):
    result = {}
    for slide_number, slide in enumerate(presentation.slides):
        result[slide_number] = ""
        # print(f"Slide {slide_number + 1}")
        slide_text = extract_text_from_slide(slide)
        slide_image_text = extract_text_from_slide_images(slide, source)
        slide_table_text = extract_text_from_slide_tables(slide)
        result[slide_number] += (
            slide_text + "\n" + slide_image_text + "\n" + slide_table_text
        )
    return result


def extract_text_from_slide_tables(slide) -> str:
    result = ""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    result += " " + cell.text
    return result


def extract_text_from_slide(slide) -> str:
    result = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    result += " " + run.text
    return result


import traceback
from tempfile import TemporaryDirectory

import pytesseract
from PIL import Image


def extract_text_from_slide_images(slide, source) -> str:
    result = ""
    for slide_shape in slide.shapes:
        try:
            if img := slide_shape._element.xpath(".//p:blipFill/a:blip/@r:embed"):
                with TemporaryDirectory() as tempdir:
                    supported_formats = {"jpeg", "png", "tiff"}
                    img_part = slide_shape.part.related_part(img[0])
                    image_bytes = img_part.blob
                    image_format = img_part.content_type.split("/")[-1]
                    if image_format not in supported_formats:
                        continue
                    image_filename = f"{tempdir}/extracted_image.{image_format}"
                    with open(image_filename, "wb") as img_file:
                        img_file.write(image_bytes)
                    # Convert image to a more stable format if necessary
                    with Image.open(image_filename) as img:
                        converted_filename = f"{tempdir}/converted_image.png"
                        img.save(converted_filename)
                    result += "\n" + str(
                        pytesseract.image_to_string(Image.open(converted_filename))
                    ).replace("\n", "")
        except Exception as e:
            supplementary_error_logger.error("%s | %s", source, str(e), exc_info=True)
    return result


if __name__ == "__main__":
    ...
